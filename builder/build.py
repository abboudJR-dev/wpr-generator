"""Orchestrator: wires extractors -> WPRState -> 15 slides -> .pptx file.

Two entry points:

* `state_from_inputs(dcrs, logs, ...)` — builds a WPRState from raw extracted data
  using the Week-20 defaults wherever editable narrative is needed.
* `build_pptx(state, output_path)` — renders the deck.
"""
from __future__ import annotations

from dataclasses import replace
from datetime import date as Date
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Inches

from extractors.dcr import DCRData, DCRPhoto
from extractors.logs import LogsData

from .defaults import (
    DEFAULT_AOCS, DEFAULT_AOC_ASSESSMENT,
    DEFAULT_CONCRETE, DEFAULT_LOOKAHEAD_BULLETS, DEFAULT_NCRS, DEFAULT_NCR_NOTE,
    DEFAULT_PROGRAMME_ROWS, DEFAULT_PROGRAMME_STATEMENT,
    DEFAULT_WEEK_SUMMARY, DEFAULT_ZONE_A_ACTIVITIES, DEFAULT_ZONE_A_LOOKAHEAD,
    DEFAULT_ZONE_B_ACTIVITIES, DEFAULT_ZONE_B_LOOKAHEAD, default_manpower_days,
)


def _default_caption_pair(dcr: DCRData) -> tuple[str, str]:
    """Return a generic caption pair for a day's photo slide. We intentionally
    do NOT use hardcoded Week-20 narrative — those captions described a
    specific construction phase and were wrong for every other week. The user
    edits captions in the Photos tab; this is just a placeholder that won't
    mislead the reader if left untouched."""
    if dcr.date:
        weekday = WEEKDAY_NAMES[dcr.date.weekday()]
        date_str = dcr.date.strftime("%d %b %Y")
        return (
            f"Site progress — {weekday}, {date_str}.",
            f"Site progress — {weekday}, {date_str}.",
        )
    return ("Site progress photograph.", "Site progress photograph.")
from .slides.activities import build_slide_current_week, build_slide_lookahead
from .slides.aoc import build_slide_aoc
from .slides.cover import build_slide_cover, build_slide_info_summary, build_slide_overview
from .slides.manpower import build_slide_manpower, default_observations
from .slides.photos import build_all_photo_slides
from .slides.programme import build_slide_programme
from .slides.quality import build_slide_quality
from .state import DayPhotos, ManpowerDay, WPRState


WEEKDAY_NAMES = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
WEEKDAY_SHORT = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]

# Project week 1 starts Mon 15 Dec 2025 — derived from the reference deck
# (Week 20 = 27 Apr → 03 May 2026 implies week 1 starts 19 weeks earlier).
# Used to auto-detect the week number from the first DCR's date.
WEEK1_START = Date(2025, 12, 15)
PROJECT_TOTAL_WEEKS = 104  # 24-month contract


def infer_week_no(first_dcr_date: Date | None) -> int:
    """Return the project week number that contains the given date."""
    if first_dcr_date is None:
        return 1
    days = (first_dcr_date - WEEK1_START).days
    return max(1, (days // 7) + 1)


def state_from_inputs(
    dcrs: list[DCRData],
    logs: LogsData | None,
    *,
    week_no: int | None = None,
    issue_date: str | None = None,
    report_ref: str | None = None,
    next_week_period: str | None = None,
    photo_picks: list[tuple[int, int]] | None = None,
    photo_captions: list[list[str]] | None = None,
) -> WPRState:
    """Assemble a WPRState ready for the slide builders.

    All meta fields auto-derive from the DCR dates if not supplied:
      * `week_no` ← project week containing the first DCR date
      * `issue_date` ← the Wednesday after the period ends (typical issue day)
      * `report_ref` ← `DB36-EPCC-WR-<NNN>` from `week_no`
      * `next_week_period` ← the seven days after the last DCR date

    `photo_picks`: list of (day_index, photo_index_within_dcr) — defaults to first
    two photos of each day.
    """
    # ----- Auto-detect week / period from DCR dates -----
    first_dt = dcrs[0].date if dcrs and dcrs[0].date else Date(2026, 4, 27)
    last_dt = dcrs[-1].date if dcrs and dcrs[-1].date else Date(2026, 5, 3)

    if week_no is None:
        week_no = infer_week_no(first_dt)
    if issue_date is None:
        # Reports are typically issued the Wednesday after the reporting period ends
        from datetime import timedelta
        issue_dt = last_dt + timedelta(days=3)
        issue_date = issue_dt.strftime("%d %B %Y")

    state = WPRState(week_no=week_no, week_no_str=str(week_no))

    state.period = f"{first_dt.strftime('%d %B')} – {last_dt.strftime('%d %B %Y')}"
    state.period_short = f"{first_dt.strftime('%d %b')} – {last_dt.strftime('%d %b %Y')}"
    state.issue_date = issue_date
    state.report_ref = report_ref or f"DB36-EPCC-WR-{str(week_no).zfill(3)}"
    state.next_week_no = week_no + 1
    state.next_week_period = next_week_period or _infer_next_week(dcrs)

    # ----- Programme progress derived from week_no -----
    state.weeks_total = PROJECT_TOTAL_WEEKS
    state.weeks_elapsed = max(1, week_no + 1)  # weeks elapsed includes the current
    pct = (state.weeks_elapsed / state.weeks_total) * 100
    state.time_elapsed_pct = f"{pct:.1f}%"

    # ----- Manpower (extracted actuals + default planned) -----
    day_labels = []
    actuals = []
    for d in dcrs[:7]:
        if d.date:
            label = f"{WEEKDAY_SHORT[d.date.weekday()]} {d.date.strftime('%d-%b')}"
        else:
            label = "Day"
        day_labels.append(label)
        actuals.append(d.subcontractor or 0)
    while len(day_labels) < 7:
        day_labels.append(f"Day {len(day_labels) + 1}")
        actuals.append(0)
    state.manpower_days = default_manpower_days(day_labels, actuals)
    state.manpower_observations = default_observations(state)

    # ----- Activities + lookahead (Week-20 defaults) -----
    state.zone_a_activities = list(DEFAULT_ZONE_A_ACTIVITIES)
    state.zone_b_activities = list(DEFAULT_ZONE_B_ACTIVITIES)
    state.week_summary_text = DEFAULT_WEEK_SUMMARY
    state.zone_a_lookahead = list(DEFAULT_ZONE_A_LOOKAHEAD)
    state.zone_b_lookahead = list(DEFAULT_ZONE_B_LOOKAHEAD)
    state.lookahead_focus_bullets = list(DEFAULT_LOOKAHEAD_BULLETS)

    # ----- NCRs + concrete -----
    state.ncrs = list(DEFAULT_NCRS)
    state.ncr_note = DEFAULT_NCR_NOTE
    state.concrete_rows = list(DEFAULT_CONCRETE)

    # ----- Submittals (from logs) -----
    if logs:
        state.submittal_categories = list(logs.categories)
        rfi = logs.by_sheet("RFI")
        if rfi:
            state.rfi_footnote = (
                f"† RFIs use R (Received) — shown in B column. "
                f"{rfi.rfi_received} received / {rfi.rfi_under_review} under review."
            )
    if not state.rfi_footnote:
        state.rfi_footnote = (
            "† RFIs use R (Received) — shown in B column. "
            "(Counts unavailable.)"
        )

    # ----- AOCs + programme -----
    state.aocs = list(DEFAULT_AOCS)
    state.aoc_assessment = DEFAULT_AOC_ASSESSMENT
    state.programme_rows = list(DEFAULT_PROGRAMME_ROWS)
    state.programme_statement = DEFAULT_PROGRAMME_STATEMENT

    # ----- Photo slides (6 days = days 0..5) -----
    state.photo_days = []
    state.photo_dates = []
    state.photo_day_labels = []
    for i, dcr in enumerate(dcrs[:6]):
        # Default: first two photos with default captions
        photo_a = dcr.photos[0] if len(dcr.photos) > 0 else None
        photo_b = dcr.photos[1] if len(dcr.photos) > 1 else None

        if photo_picks and i < len(photo_picks):
            a_idx, b_idx = photo_picks[i]
            photo_a = dcr.photos[a_idx] if 0 <= a_idx < len(dcr.photos) else photo_a
            photo_b = dcr.photos[b_idx] if 0 <= b_idx < len(dcr.photos) else photo_b

        if photo_captions and i < len(photo_captions) and len(photo_captions[i]) >= 2:
            cap_a, cap_b = photo_captions[i][0], photo_captions[i][1]
        else:
            cap_a, cap_b = _default_caption_pair(dcr)

        state.photo_days.append(DayPhotos(
            photo_a_bytes=photo_a.bytes_ if photo_a else None,
            photo_a_ext=photo_a.ext if photo_a else "png",
            photo_b_bytes=photo_b.bytes_ if photo_b else None,
            photo_b_ext=photo_b.ext if photo_b else "png",
            caption_a=cap_a,
            caption_b=cap_b,
        ))
        if dcr.date:
            state.photo_dates.append(dcr.date)
            state.photo_day_labels.append(WEEKDAY_NAMES[dcr.date.weekday()])
        else:
            state.photo_dates.append(Date(2026, 4, 27 + i))
            state.photo_day_labels.append(WEEKDAY_NAMES[i % 7])

    return state


def _infer_next_week(dcrs: list[DCRData]) -> str:
    if not dcrs or not dcrs[-1].date:
        return "04 – 10 May 2026"
    last = dcrs[-1].date
    next_start = Date.fromordinal(last.toordinal() + 1)
    next_end = Date.fromordinal(last.toordinal() + 7)
    return f"{next_start.strftime('%d')} – {next_end.strftime('%d %B %Y')}"


def build_pptx(state: WPRState, output_path: str | Path) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_cover(prs, state)             # 1
    build_slide_overview(prs, state)          # 2
    build_slide_info_summary(prs, state)      # 3
    build_slide_manpower(prs, state)          # 4
    build_slide_current_week(prs, state)      # 5
    build_slide_lookahead(prs, state)         # 6
    build_slide_quality(prs, state)           # 7
    build_slide_aoc(prs, state)               # 8
    build_slide_programme(prs, state)         # 9
    build_all_photo_slides(prs, state, start_page=9)  # 10-15

    prs.save(str(output_path))
    state.output_path = output_path
    return output_path
