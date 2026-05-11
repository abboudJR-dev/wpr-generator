"""Slide 7: Quality & NCR Status (NCR register, concrete quality, submittal dashboard, HSE)."""
from __future__ import annotations

from pptx.util import Inches, Pt

from builder.design import (
    C, Run, SLIDE_W,
    add_blank_slide, add_footer, add_header, add_oval, add_rect, add_text_box,
    fill_cell, kpi_card, section_label, set_cell_border, set_slide_background,
    style_cell_text,
)
from builder.state import ConcreteRow, NCREntry, WPRState
from extractors.logs import CategoryStats


def _add_ncr_table(slide, x, y, w, ncrs: list[NCREntry]):
    n_rows = len(ncrs) + 1
    n_cols = 5
    col_widths = [0.85, 2.95, 0.85, 0.85, 1.00]
    row_h = 0.28

    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                         Inches(w), Inches(row_h * n_rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)
    for ri in range(n_rows):
        table.rows[ri].height = Inches(row_h)

    headers = ["NCR REF.", "SUBJECT", "ISSUED", "CLOSED", "STATUS"]
    aligns = ["center", "left", "center", "center", "center"]
    for ci, (h, a) in enumerate(zip(headers, aligns)):
        cell = table.cell(0, ci)
        fill_cell(cell, C.navy)
        style_cell_text(cell, h, color=C.white, bold=True, size=8.5,
                        align=a, valign="middle", char_spacing=1.5)
        set_cell_border(cell, C.line, 0.5)

    for ri, n in enumerate(ncrs, start=1):
        stripe = C.white if ri % 2 == 1 else C.panel
        row_data = [
            (n.ref,     C.ink,      True,  "center"),
            (n.subject, C.graphite, False, "left"),
            (n.issued,  C.graphite, False, "center"),
            (n.closed,  C.graphite, False, "center"),
            (n.status,  C.success,  True,  "center"),
        ]
        for ci, (text, color, bold, align) in enumerate(row_data):
            cell = table.cell(ri, ci)
            fill_cell(cell, stripe)
            style_cell_text(cell, text, color=color, bold=bold, size=8.5,
                            align=align, valign="middle")
            set_cell_border(cell, C.line, 0.5)


def _add_concrete_table(slide, x, y, w, rows: list[ConcreteRow]):
    n_rows = len(rows) + 1
    n_cols = 5
    col_widths = [1.95, 0.75, 1.30, 1.30, 1.20]
    row_h = 0.30

    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                         Inches(w), Inches(row_h * n_rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)
    for ri in range(n_rows):
        table.rows[ri].height = Inches(row_h)

    headers = ["ELEMENT", "GRADE", "3-DAY", "7-DAY", "28-DAY"]
    aligns = ["left", "center", "center", "center", "center"]
    for ci, (h, a) in enumerate(zip(headers, aligns)):
        cell = table.cell(0, ci)
        fill_cell(cell, C.navy)
        style_cell_text(cell, h, color=C.white, bold=True, size=8.5,
                        align=a, valign="middle", char_spacing=1.5)
        set_cell_border(cell, C.line, 0.5)

    for ri, row in enumerate(rows, start=1):
        stripe = C.white if ri % 2 == 1 else C.panel
        stat28_color = C.amber if row.day28_kind == "pending" else C.success
        cells_def = [
            (row.element, C.ink,      True,  "left"),
            (row.grade,   C.graphite, False, "center"),
            (row.day3,    C.success,  False, "center"),
            (row.day7,    C.success,  False, "center"),
            (row.day28,   stat28_color, True, "center"),
        ]
        for ci, (text, color, bold, align) in enumerate(cells_def):
            cell = table.cell(ri, ci)
            fill_cell(cell, stripe)
            style_cell_text(cell, text, color=color, bold=bold, size=8.5,
                            align=align, valign="middle")
            set_cell_border(cell, C.line, 0.5)


def _add_submittal_table(slide, x, y, w, categories: list[CategoryStats]):
    n_rows = len(categories) + 1
    n_cols = 7
    col_widths = [1.95, 0.70, 0.55, 0.55, 0.55, 0.55, 0.55]
    row_h = 0.28

    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                         Inches(w), Inches(row_h * n_rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)
    for ri in range(n_rows):
        table.rows[ri].height = Inches(row_h)

    # Header row (column-coloured for A/B/C-D/E/F)
    header_specs = [
        ("CATEGORY",  C.navy,    "left",   1),
        ("TOTAL",     C.navy,    "center", 1),
        ("A",         C.success, "center", 0),
        ("B",         C.teal,    "center", 0),
        ("C/D",       C.amber,   "center", 0),
        ("E",         C.slate,   "center", 0),
        ("F",         C.steel,   "center", 0),
    ]
    for ci, (text, fill, align, char_spacing) in enumerate(header_specs):
        cell = table.cell(0, ci)
        fill_cell(cell, fill)
        style_cell_text(cell, text, color=C.white, bold=True, size=8.5,
                        align=align, valign="middle", char_spacing=char_spacing)
        set_cell_border(cell, C.line, 0.5)

    # Data rows
    col_text_colors = [None, C.graphite, C.success, C.teal_dark, C.amber, C.slate, C.steel]
    for ri, cat in enumerate(categories, start=1):
        stripe = C.white if ri % 2 == 1 else C.panel
        row_vals = cat.display_row()
        for ci, val in enumerate(row_vals):
            cell = table.cell(ri, ci)
            fill_cell(cell, stripe)
            if ci == 0:
                style_cell_text(cell, val, color=C.ink, bold=False, size=8.5,
                                align="left", valign="middle")
            elif ci == 1:
                style_cell_text(cell, val, color=C.graphite, bold=True, size=8.5,
                                align="center", valign="middle")
            else:
                color = col_text_colors[ci]
                bold = val not in ("—", "0")
                style_cell_text(cell, val, color=color, bold=bold, size=8.5,
                                align="center", valign="middle")
            set_cell_border(cell, C.line, 0.5)


def build_slide_quality(prs, state: WPRState):
    slide = add_blank_slide(prs)
    set_slide_background(slide, C.page)
    # Subtitle uses the period end date (not the issue date), matching the reference
    as_at = state.period_as_at()
    add_header(
        slide,
        title="QUALITY & NCR STATUS SUMMARY",
        subtitle=f"Week {state.week_no_str}  |  As at {as_at}",
    )
    add_footer(slide, 6, period=state.period_short,
               company=state.project.contractor_short,
               report_no=f"Weekly Report No. {state.week_no_str.zfill(3)}")

    # LEFT: NCR register
    section_label(slide, 0.45, 1.05, 6.5, "NON-CONFORMANCE REPORT (NCR) REGISTER")
    _add_ncr_table(slide, 0.45, 1.40, 6.50, state.ncrs)

    # NCR commentary callout
    ncr_note_y = 3.50
    add_rect(slide, 0.45, ncr_note_y, 6.50, 0.65, C.success_soft)
    add_rect(slide, 0.45, ncr_note_y, 0.06, 0.65, C.success)
    closed = sum(1 for n in state.ncrs if "CLOSED" in n.status.upper())
    add_text_box(
        slide, 0.65, ncr_note_y, 6.20, 0.65,
        [
            Run(text=f"All {closed} NCRs closed. ", bold=True, color=C.success, size=9),
            Run(text=state.ncr_note, color=C.graphite, size=9),
        ],
        valign="middle", margin=0,
    )

    # Concrete Quality Summary
    section_label(slide, 0.45, 4.30, 6.5, "CONCRETE QUALITY SUMMARY")
    _add_concrete_table(slide, 0.45, 4.65, 6.50, state.concrete_rows)

    # RIGHT: Submittal Status Dashboard
    section_label(slide, 7.15, 1.05, 5.7, "SUBMITTAL STATUS DASHBOARD")
    _add_submittal_table(slide, 7.15, 1.40, 5.75, state.submittal_categories)

    # Legend
    lg_y = 4.40
    legend_items = [
        (C.success, "A — Approved"),
        (C.teal,    "B — Approved as Noted"),
        (C.amber,   "C/D — Revise & Resubmit / Rejected"),
        (C.slate,   "E — Under Review"),
        (C.steel,   "F — For Information Only"),
    ]
    for i, (color, label) in enumerate(legend_items):
        lx = 7.15 + (i % 2) * 2.85
        ly = lg_y + (i // 2) * 0.26
        add_rect(slide, lx, ly + 0.05, 0.16, 0.16, color)
        add_text_box(
            slide, lx + 0.22, ly, 2.6, 0.26,
            [Run(text=label, size=8, color=C.graphite)],
            valign="middle", margin=0,
        )

    # RFI footnote
    add_text_box(
        slide, 7.15, lg_y + 0.85, 5.75, 0.22,
        [Run(text=state.rfi_footnote, size=7.5, italic=True, color=C.slate)],
        valign="middle", margin=0,
    )

    # BOTTOM: HSE Weekly Summary
    section_label(slide, 0.45, 5.65, 12.5, "HSE WEEKLY SUMMARY")
    hse_y = 5.95
    hse = state.hse
    cards = [
        (str(hse.lti),         "Lost Time Injuries (LTI)",     C.success, None),
        (str(hse.open_ncrs),   "Open NCRs",                    C.success, None),
        (str(hse.closed_ncrs), "NCRs Closed to Date",          C.teal,    None),
        ("✓",                  f"Safety Officer CV ({hse.safety_officer_doc})",
                               C.steel, hse.safety_officer_status),
    ]
    hse_w = (SLIDE_W - 0.90 - 0.30) / 4
    for i, (v, label, fill, val_label) in enumerate(cards):
        cx = 0.45 + i * (hse_w + 0.10)
        add_rect(slide, cx, hse_y, hse_w, 1.00, fill)
        if val_label:
            add_text_box(
                slide, cx, hse_y + 0.15, hse_w, 0.45,
                [Run(text=val_label, size=20, bold=True, color=C.white, char_spacing=2)],
                align="center", valign="middle", margin=0,
            )
        else:
            add_text_box(
                slide, cx, hse_y + 0.05, hse_w, 0.60,
                [Run(text=v, size=32, bold=True, color=C.white)],
                align="center", valign="middle", margin=0,
            )
        add_text_box(
            slide, cx, hse_y + 0.70, hse_w, 0.26,
            [Run(text=label.upper(), size=8.5, color=C.white, char_spacing=2)],
            align="center", valign="middle", margin=0,
        )

    return slide


# Helper added on the state class via monkey-patch — keeps state.py minimal.
def _period_as_at(self) -> str:
    """Format the period end-date as e.g. '03 May 2026' for slide subtitles."""
    if self.photo_dates:
        from datetime import timedelta
        end_dt = self.photo_dates[-1] + timedelta(days=1)
        return end_dt.strftime("%d %B %Y")
    return self.issue_date


WPRState.period_as_at = _period_as_at  # type: ignore[attr-defined]
