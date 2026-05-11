"""Design tokens + reusable layout primitives for the WPR deck.

All coordinates are in inches; the slide is 13.333" x 7.5" (LAYOUT_WIDE 16:9).
This is a 1:1 port of `reference/scaffolding/build.js` from pptxgenjs to python-pptx.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Optional, Sequence

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

import os
import sys

# ----- Slide geometry -----
SLIDE_W = 13.333
SLIDE_H = 7.5


def _project_root() -> Path:
    """Where the bundled assets and source live, both in dev and inside a
    PyInstaller bundle (--onedir or --onefile)."""
    if getattr(sys, "frozen", False):
        if hasattr(sys, "_MEIPASS"):
            return Path(sys._MEIPASS)  # type: ignore[attr-defined]
        return Path(sys.executable).parent
    return Path(__file__).resolve().parent.parent


ROOT = _project_root()
ASSETS = {
    "logo_ew":   ROOT / "assets" / "logo_eastwest.png",
    "logo_ah":   ROOT / "assets" / "logo_alhadara.png",
    "logo_epc":  ROOT / "assets" / "logo_emiratespearl.png",
    "building":  ROOT / "assets" / "building_render.jpg",
}


# ----- Color palette (from build.js C object) -----
@dataclass(frozen=True)
class Palette:
    navy: str        = "0F2A47"
    navy_deep: str   = "0A1F36"
    steel: str       = "1E5A8C"
    teal: str        = "0E9F8F"
    teal_dark: str   = "0C8074"
    teal_soft: str   = "E5F4F2"
    success: str     = "1F7A4D"
    success_soft: str = "E3F1E9"
    amber: str       = "C8881B"
    amber_soft: str  = "FCEFD8"
    risk: str        = "B33A2A"
    risk_soft: str   = "F8E1DC"
    pending: str     = "5A6B7D"
    pending_soft: str = "EAEEF2"
    ink: str         = "1A2332"
    graphite: str    = "3A4A5C"
    slate: str       = "5A6B7D"
    mute: str        = "8A99A8"
    line: str        = "D8DEE5"
    hairline: str    = "ECEFF3"
    page: str        = "FFFFFF"
    panel: str       = "F5F7FA"
    panel_alt: str   = "EFF3F8"
    white: str       = "FFFFFF"
    cover_strip: str = "07172B"


C = Palette()


def hex_to_rgb(hex_str: str) -> RGBColor:
    """'0F2A47' -> RGBColor(0x0F, 0x2A, 0x47)."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ----- Low-level shape / text builders ----------------------------------------

def add_rect(slide, x, y, w, h, fill_hex: str, line_hex: Optional[str] = None,
             line_width: Optional[float] = None):
    """Solid rectangle. line_hex defaults to fill (no visible border)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex is None:
        shape.line.color.rgb = hex_to_rgb(fill_hex)
    else:
        shape.line.color.rgb = hex_to_rgb(line_hex)
    if line_width is not None:
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_oval(slide, x, y, w, h, fill_hex: str, line_hex: Optional[str] = None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shape.line.color.rgb = hex_to_rgb(line_hex or fill_hex)
    shape.shadow.inherit = False
    return shape


def add_image(slide, path: str | Path, x, y, w, h):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


# Text: a "run" specifies one styled chunk, paragraphs are joined with breaks.
@dataclass
class Run:
    text: str
    bold: bool = False
    italic: bool = False
    color: str = C.ink
    size: float = 10.0
    font: str = "Calibri"
    char_spacing: float = 0.0
    underline: bool = False
    new_paragraph: bool = False  # if True, start a new paragraph BEFORE this run


def _apply_run_format(run, run_def: Run):
    run.text = run_def.text
    run.font.name = run_def.font
    run.font.size = Pt(run_def.size)
    run.font.bold = run_def.bold
    run.font.italic = run_def.italic
    run.font.color.rgb = hex_to_rgb(run_def.color)
    if run_def.underline:
        run.font.underline = True
    if run_def.char_spacing:
        # python-pptx doesn't expose 'spc' (kerning) directly — set via XML
        rPr = run._r.get_or_add_rPr()
        # spc is in 1/100 of a point; +N adds N hundredths-of-pt of letter spacing
        rPr.set("spc", str(int(run_def.char_spacing * 100)))


def add_text_box(slide, x, y, w, h, runs: Sequence[Run] | str, *,
                 align: str = "left", valign: str = "top", margin: float = 0.0,
                 word_wrap: bool = True):
    """Place a text frame at (x,y,w,h). `runs` is either a plain string (single run)
    or a sequence of Run instances.

    align: 'left' | 'center' | 'right'
    valign: 'top' | 'middle' | 'bottom'
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = word_wrap
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[valign]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    if isinstance(runs, str):
        runs = [Run(text=runs)]

    # Auto-split runs that contain '\n' into multiple sub-runs (one per paragraph).
    expanded: list[Run] = []
    for r in runs:
        if "\n" in r.text:
            parts = r.text.split("\n")
            for idx, part in enumerate(parts):
                expanded.append(Run(
                    text=part, bold=r.bold, italic=r.italic, color=r.color,
                    size=r.size, font=r.font, char_spacing=r.char_spacing,
                    underline=r.underline,
                    new_paragraph=(r.new_paragraph if idx == 0 else True),
                ))
        else:
            expanded.append(r)

    p = tf.paragraphs[0]
    p.alignment = align_map[align]
    first_run_in_paragraph = True

    for r in expanded:
        if r.new_paragraph and not first_run_in_paragraph:
            p = tf.add_paragraph()
            p.alignment = align_map[align]
            first_run_in_paragraph = True
        run = p.add_run()
        _apply_run_format(run, r)
        first_run_in_paragraph = False

    return box


# ----- Reusable composite primitives ------------------------------------------

def add_hairline(slide, y: float, color: str, h: float = 0.06):
    add_rect(slide, 0, y, SLIDE_W, h, color)


def add_header(slide, title: str, subtitle: Optional[str] = None):
    """White header band: navy hairline top (0.06"), teal hairline bottom (0.02"),
    title left, three logos right."""
    add_rect(slide, 0, 0, SLIDE_W, 0.06, C.navy)              # top hairline
    add_rect(slide, 0, 0.06, SLIDE_W, 0.84, C.page)            # white band
    add_rect(slide, 0, 0.90, SLIDE_W, 0.02, C.teal)            # teal hairline

    add_text_box(
        slide, 0.45, 0.13, 8.0, 0.40,
        [Run(text=title, size=17, bold=True, color=C.navy)],
        align="left", valign="middle", margin=0,
    )
    if subtitle:
        add_text_box(
            slide, 0.45, 0.52, 8.0, 0.30,
            [Run(text=subtitle, size=9.5, color=C.slate)],
            align="left", valign="middle", margin=0,
        )

    # Logos top-right
    logo_y = 0.16
    logo_h = 0.58
    add_image(slide, ASSETS["logo_ew"],  9.55,  logo_y,        1.15, logo_h)
    add_image(slide, ASSETS["logo_ah"],  10.85, logo_y + 0.08, 1.15, logo_h - 0.16)
    add_image(slide, ASSETS["logo_epc"], 12.15, logo_y,        0.85, logo_h)


def add_footer(slide, page_num: int, *,
               company: str = "Emirates Pearl Construction CCS Company",
               report_no: str = "Weekly Report No. 020",
               period: str = "27 Apr – 03 May 2026"):
    add_rect(slide, 0, SLIDE_H - 0.34, SLIDE_W, 0.34, C.navy)

    sep = Run(text="  |  ", color="8FA8C2", size=8.5)
    runs = [
        Run(text=company, bold=True, color=C.white, size=8.5),
        sep,
        Run(text=report_no, color="DCE6F2", size=8.5),
        sep,
        Run(text=period, color="DCE6F2", size=8.5),
        sep,
        Run(text="CONFIDENTIAL", bold=True, color=C.teal, size=8.5, char_spacing=2),
    ]
    add_text_box(
        slide, 0.45, SLIDE_H - 0.34, 11.5, 0.34,
        runs, align="left", valign="middle", margin=0,
    )
    add_text_box(
        slide, SLIDE_W - 0.85, SLIDE_H - 0.34, 0.45, 0.34,
        [Run(text=str(page_num).zfill(2), bold=True, color=C.white, size=9.5)],
        align="right", valign="middle", margin=0,
    )


def section_label(slide, x: float, y: float, w: float, text: str, color: str = C.steel):
    add_text_box(
        slide, x, y, w, 0.28,
        [Run(text=text, size=10.5, bold=True, color=color, char_spacing=4)],
        align="left", valign="middle", margin=0,
    )


def zone_banner(slide, x: float, y: float, w: float, label: str, color: str):
    add_rect(slide, x, y, w, 0.32, color)
    add_text_box(
        slide, x + 0.15, y, w - 0.30, 0.32,
        [Run(text=label, size=10.5, bold=True, color=C.white, char_spacing=1.5)],
        align="left", valign="middle", margin=0,
    )


def status_pill(slide, x: float, y: float, w: float, h: float, label: str, kind: str = "info"):
    cfg_map = {
        "completed": {"fill": C.success_soft, "text": C.success,   "label": "COMPLETED"},
        "progress":  {"fill": C.teal_soft,    "text": C.teal_dark, "label": "IN PROGRESS"},
        "pending":   {"fill": C.pending_soft, "text": C.pending,   "label": "PENDING"},
        "risk":      {"fill": C.risk_soft,    "text": C.risk,      "label": "AT RISK"},
        "info":      {"fill": C.panel_alt,    "text": C.steel,     "label": label or "INFO"},
    }
    cfg = cfg_map.get(kind, cfg_map["info"])
    add_rect(slide, x, y, w, h, cfg["fill"])
    add_text_box(
        slide, x, y, w, h,
        [Run(text=label or cfg["label"], size=8.5, bold=True, color=cfg["text"], char_spacing=1)],
        align="center", valign="middle", margin=0,
    )


def kpi_card(slide, x: float, y: float, w: float, h: float, *,
             value: str, label: str, fill: str, value_size: float = 22,
             value_color: str = C.white, label_size: float = 8.5,
             label_color: str = C.white, label_char_spacing: float = 2.5):
    """Solid-color card with big number top + small label bottom (centered)."""
    add_rect(slide, x, y, w, h, fill)
    add_text_box(
        slide, x, y + 0.10, w, h * 0.55,
        [Run(text=value, size=value_size, bold=True, color=value_color)],
        align="center", valign="middle", margin=0,
    )
    add_text_box(
        slide, x, y + h - 0.30, w, 0.25,
        [Run(text=label.upper(), size=label_size, color=label_color, char_spacing=label_char_spacing)],
        align="center", valign="middle", margin=0,
    )


def kpi_card_inline(slide, x: float, y: float, w: float, h: float, *,
                    value: str, label: str, fill: str,
                    value_size: float = 32, label_size: float = 9):
    """Variant: big number on the left, label stacked on the right."""
    add_rect(slide, x, y, w, h, fill)
    add_text_box(
        slide, x + 0.20, y + 0.08, 1.4, h - 0.20,
        [Run(text=value, size=value_size, bold=True, color=C.white)],
        align="left", valign="middle", margin=0,
    )
    add_text_box(
        slide, x + 1.65, y + 0.30, w - 1.85, 0.40,
        [Run(text=label.upper(), size=label_size, bold=True, color=C.white, char_spacing=2)],
        align="left", valign="middle", margin=0,
    )


# ----- Table cell formatting helpers ------------------------------------------

def fill_cell(cell, hex_color: str):
    cell.fill.solid()
    cell.fill.fore_color.rgb = hex_to_rgb(hex_color)


def style_cell_text(cell, text: str, *, color: str = C.ink, bold: bool = False,
                    size: float = 9.0, align: str = "left", valign: str = "middle",
                    char_spacing: float = 0.0, font: str = "Calibri"):
    """Replace the cell text with a single styled run."""
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM
    }[valign]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    tf = cell.text_frame
    tf.word_wrap = True
    # Reset to one paragraph + one run
    p = tf.paragraphs[0]
    p.alignment = align_map[align]
    # Clear any existing runs
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    # If paragraph still has text from before, blank it via text setter
    if p.text:
        p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    if char_spacing:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(char_spacing * 100)))


def set_cell_border(cell, color: str = C.line, width_pt: float = 0.5):
    """Add 4-side border to a table cell via XML (python-pptx has no public API for this)."""
    tcPr = cell._tc.get_or_add_tcPr()
    color_clean = color.lstrip("#")
    width_emu = Pt(width_pt)
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for border_tag in ("lnL", "lnR", "lnT", "lnB"):
        # Remove existing
        for existing in tcPr.findall(qn(f"a:{border_tag}")):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(f"a:{border_tag}"))
        ln.set("w", str(int(width_emu)))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        fill = etree.SubElement(ln, qn("a:solidFill"))
        clr = etree.SubElement(fill, qn("a:srgbClr"))
        clr.set("val", color_clean)
        prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", "solid")


def add_blank_slide(prs):
    """Add a blank slide (no placeholders) and return it."""
    blank_layout = prs.slide_layouts[6]  # 'Blank' layout in default templates
    return prs.slides.add_slide(blank_layout)


def add_bulleted_text(slide, x, y, w, h, items, *,
                      bullet_char: str = "■", color: str = C.graphite,
                      size: float = 9.5, font: str = "Calibri",
                      bullet_indent: float = 0.20):
    """Bulleted list using paragraph-level bullet character (a:buChar)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for idx, text in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Clear default runs
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        # Add bullet character via pPr/buChar
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        # Remove existing bullets
        for child in list(pPr):
            if child.tag in (qn("a:buChar"), qn("a:buAutoNum"), qn("a:buNone")):
                pPr.remove(child)
        # marL + indent for bullet alignment
        pPr.set("marL", str(int(bullet_indent * 914400)))
        pPr.set("indent", str(int(-bullet_indent * 914400)))
        bu_char = etree.SubElement(pPr, qn("a:buChar"))
        bu_char.set("char", bullet_char)
        # Run with the body text
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = hex_to_rgb(color)

    return box


def set_slide_background(slide, hex_color: str):
    """Solid-color slide background via XML."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)
